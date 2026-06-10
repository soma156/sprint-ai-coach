"""
将 AR眼镜第一人称篮球观赛体验.pptx 转换为华南师范大学PPT模板样式
v3 - 使用安全的XML操作方式
"""
import os
import re
import tempfile
import zipfile
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ============ SCNU Brand Colors ============
SCNU_RED = '8B1A2B'
SCNU_DARK = '2D2D2D'
SCNU_WHITE = 'FFFFFF'
SCNU_GRAY = '666666'
SCNU_LIGHT_GRAY = '888888'
SCNU_CARD_BG = 'F7F3F3'
SCNU_CARD_BG2 = 'FAF8F7'

# Color mappings
FILL_MAP = {
    'FF6B2B': SCNU_RED,   # orange -> red
    '112B50': SCNU_CARD_BG,  # dark navy -> light card
    '0B1D3A': SCNU_CARD_BG2, # darker navy -> lighter card
    '000000': SCNU_RED,   # black -> red
    '4A90D9': SCNU_RED,   # blue -> red
}

TEXT_MAP = {
    'FFFFFF': SCNU_DARK,  # white -> dark
    'A0AEC0': SCNU_GRAY,  # light gray -> gray
    '4A90D9': SCNU_RED,   # blue -> red
    'FF6B2B': SCNU_RED,   # orange -> red
}

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

print("=" * 60)
print("华南师范大学PPT模板样式转换 v3")
print("=" * 60)

# ============ Step 1: Extract logo ============
print("\n[1/4] 提取华师Logo...")
prs_scnu = Presentation('AR篮球POV_华师模板.pptx')
logo_blob = None
logo_content_type = None
for sh in prs_scnu.slides[0].shapes:
    if sh.name == '图片 5':
        logo_blob = sh.image.blob
        logo_content_type = sh.image.content_type
        print(f"  Logo: {len(logo_blob)} bytes ({logo_content_type})")

# ============ Step 2: Process ============
print("\n[2/4] 转换颜色...")

prs = Presentation('AR眼镜第一人称篮球观赛体验.pptx')

def replace_srgb_in_element(elem, color_map):
    """Replace srgbClr values in an element tree (in-place)"""
    count = 0
    for srgb in elem.iter(f'{{{A_NS}}}srgbClr'):
        old_val = srgb.get('val')
        if old_val and old_val.upper() in color_map:
            srgb.set('val', color_map[old_val.upper()])
            count += 1
    return count

def make_sp_element(id_num, name, left, top, width, height, fill_color=None, text=None, font_size='900'):
    """Create a simple shape (rectangle or textbox) XML element"""
    if text is not None:
        # Textbox
        txBody = f'''<p:txBody xmlns:a="{A_NS}">
          <a:bodyPr/>
          <a:lstStyle/>
          <a:p>
            <a:pPr algn="r"/>
            <a:r>
              <a:rPr sz="{font_size}">
                <a:solidFill><a:srgbClr val="{SCNU_LIGHT_GRAY}"/></a:solidFill>
                <a:latin typeface="Arial"/>
                <a:ea typeface="Microsoft YaHei"/>
              </a:rPr>
              <a:t>{text}</a:t>
            </a:r>
          </a:p>
        </p:txBody>'''
    else:
        txBody = '<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'

    fill_xml = f'<a:solidFill><a:srgbClr val="{fill_color}"/></a:solidFill>' if fill_color else '<a:noFill/>'

    xml = f'''<p:sp xmlns:p="{P_NS}" xmlns:a="{A_NS}">
      <p:nvSpPr>
        <p:cNvPr id="{id_num}" name="{name}"/>
        <p:cNvSpPr/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="{left}" y="{top}"/>
          <a:ext cx="{width}" cy="{height}"/>
        </a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        {fill_xml}
        <a:ln><a:noFill/></a:ln>
      </p:spPr>
      {txBody}
    </p:sp>'''
    return etree.fromstring(xml.encode('utf-8'))

fill_stats = {k: 0 for k in FILL_MAP}
text_stats = {k: 0 for k in TEXT_MAP}

for slide_idx, slide in enumerate(prs.slides):
    spTree = slide._element  # cSld element
    # Find the actual spTree
    # slide._element is <p:sld>, spTree is <p:cSld><p:spTree>
    ns = {'p': P_NS, 'a': A_NS}
    cSld = slide._element.find(f'{{{P_NS}}}cSld')
    if cSld is None:
        print(f"  Slide {slide_idx+1}: No cSld found, skipping")
        continue
    spTree = cSld.find(f'{{{P_NS}}}spTree')
    if spTree is None:
        print(f"  Slide {slide_idx+1}: No spTree found, skipping")
        continue

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    print(f"  Slide {slide_idx+1}/10...")

    # Change background to white
    try:
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    except Exception as e:
        print(f"    BG warning: {e}")

    # Process shapes in spTree
    shape_count = 0
    for shape_elem in list(spTree):
        # Only process shape elements (p:sp, p:pic, p:grpSp, p:graphicFrame)
        tag = etree.QName(shape_elem).localname if isinstance(shape_elem, etree._Element) else None
        if tag not in ('sp', 'pic', 'grpSp', 'graphicFrame'):
            continue
        shape_count += 1
        f = replace_srgb_in_element(shape_elem, FILL_MAP)
        t = replace_srgb_in_element(shape_elem, TEXT_MAP)
        for k, v in FILL_MAP.items():
            fill_stats[k] += (1 if f > 0 and any(srgb.get('val') == v for srgb in shape_elem.iter(f'{{{A_NS}}}srgbClr')) else 0)

    # Add top red bar
    bar = make_sp_element(9000+slide_idx, 'SCNU_TopBar', 0, 0, slide_width, 54864, fill_color=SCNU_RED)
    spTree.append(bar)

    # Add bottom red line
    line_y = slide_height - 411480
    line_w = slide_width - 914400
    bline = make_sp_element(8000+slide_idx, 'SCNU_BottomLine', 457200, line_y, line_w, 18288, fill_color=SCNU_RED)
    spTree.append(bline)

    # Add page number
    pn_x = slide_width - 1371600
    pn_y = line_y + 91440
    pn = make_sp_element(7000+slide_idx, 'SCNU_PageNum', pn_x, pn_y, 914400, 274320, text=str(slide_idx+1))
    spTree.append(pn)

    print(f"    {shape_count} shapes processed")

# Print stats
print(f"\n  Fill changes: {sum(fill_stats.values())}")
for k, v in fill_stats.items():
    if v > 0: print(f"    #{k} -> #{FILL_MAP[k]}: {v}")

# ============ Step 3: Add Logo to first slide ============
print("\n[3/4] 添加Logo和装饰元素...")
if logo_blob:
    try:
        slide1 = prs.slides[0]
        ext = 'png' if 'png' in logo_content_type else 'jpg'
        with tempfile.NamedTemporaryFile(suffix=f'.{ext}', delete=False) as f:
            f.write(logo_blob)
            tmp_path = f.name
        slide1.shapes.add_picture(tmp_path, Emu(457200), Emu(640080), Emu(914400), Emu(1219200))
        os.unlink(tmp_path)
        print("  已添加华师Logo到标题页")
    except Exception as e:
        print(f"  Logo warning: {e}")

# ============ Step 4: Save ============
output_path = 'AR眼镜第一人称篮球观赛体验_华师模板.pptx'
print(f"\n[4/4] 保存到 {output_path}...")
prs.save(output_path)

# ============ Post-process theme ============
print("  更新主题配色...")
tmp_path = output_path + '.tmp'

with zipfile.ZipFile(output_path, 'r') as zin:
    with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if 'theme' in item.filename.lower() and item.filename.endswith('.xml'):
                root = etree.fromstring(data)
                theme_colors = {
                    'dk1': '000000', 'lt1': 'FFFFFF', 'dk2': SCNU_DARK, 'lt2': SCNU_CARD_BG,
                    'accent1': SCNU_RED, 'accent2': 'B33A4B', 'accent3': 'C49A2A',
                    'accent4': SCNU_RED, 'accent5': 'D4556A', 'accent6': '2D5A6B',
                    'hlink': SCNU_RED, 'folHlink': '954F72',
                }
                for cname, cval in theme_colors.items():
                    for elem in root.iter(f'{{{A_NS}}}{cname}'):
                        for srgb in elem.iter(f'{{{A_NS}}}srgbClr'):
                            srgb.set('val', cval)
                data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)
            zout.writestr(item, data)

os.replace(tmp_path, output_path)
print("  完成!")

# Verify
print(f"\n✅ 转换完成！")
print(f"📄 输出: {os.path.abspath(output_path)}")
prs_v = Presentation(output_path)
print(f"📊 {len(prs_v.slides)} 页幻灯片")
