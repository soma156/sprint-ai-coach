# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from copy import deepcopy
import os

TEMPLATE = r"D:\BaiduNetdiskDownload\华南师范大学总结汇报主题PPT模板.pptx"
OUTPUT = r"C:\Users\Aa173\新建文件夹\ai训练计划\AR篮球POV_华师模板.pptx"

prs = Presentation(TEMPLATE)

# 获取所有幻灯片
slides = list(prs.slides)

# 定义 AR POV 内容
content = {
    0: {  # 幻灯片1: 标题页 (0-indexed)
        "texts": [
            ("汇报人：", "AR眼镜第一人称篮球观赛体验"),
            ("汇报时间：", "2026.06  |  AR POV Team"),
        ],
        "special": "title_page"
    },
    1: {  # 幻灯片2: 目录
        "texts": [
            ("请输入标题", "什么是AR眼镜第一人称观赛"),
            ("请输入标题", "技术可行性：已进入实战阶段"),
            ("请输入标题", "技术架构与流程"),
            ("请输入标题", "核心优势与创新亮点"),
        ],
        "title": "目 录"
    },
    2: {  # 幻灯片3
        "title": "什么是AR眼镜第一人称观赛？",
        "body": "核心概念：通过AR眼镜，现场观众可实时切换到喜欢的篮球运动员的第一人称视角（POV），在真实赛场视野上叠加或切换球员主观画面。\n\n• 不是单纯录像，而是AI实时重建+叠加\n• 支持个性化选择（LeBron / Curry / 本地球星）\n• 结合真实环境：既能看到现场，又能代入球员视角"
    },
    3: {  # 幻灯片4
        "title": "技术可行性：已进入实战阶段",
        "body": "• NBA 2026 All-Star Technology Summit 正式预览 POV Mode\n• AI驱动实时球员视角观看\n• 核心技术：多机位拍摄 + AI 3D数字孪生球场 + 球员追踪系统\n• AR眼镜载体：Xreal、Ray-Ban Meta、Snap Spectacles\n• 已从VIP试点向常规赛事扩展"
    },
    4: {  # 幻灯片5
        "title": "技术架构与流程",
        "body": "01 场馆部署：5G + 边缘计算 + 多角度高精度摄像头\n02 数据层：Hawk-Eye追踪系统实时生成3D模型\n03 呈现层：AR眼镜接收低延迟流，AI渲染球员POV\n04 用户交互：语音/手势/App选择球员，即时切换\n\n实施路径：短期VIP租赁 → 中期全场覆盖 → 长期Metaverse远程同步"
    },
    5: {  # 幻灯片6
        "title": "实施要点与合作伙伴",
        "body": "关键要素：\n• 低延迟保障（<50ms）\n• 信息分层显示（避免过载）\n• 隐私保护（AI重建为主）\n\n推荐合作方：\n• NBA官方 / Meta / Xreal / Quintar\n• Stats Perform / Hawk-Eye\n• 5G运营商（场馆覆盖）\n\n成本模式：眼镜租赁 + 付费高级POV + 品牌AR植入赞助"
    },
    6: {  # 幻灯片7
        "title": "核心优势",
        "body": "对粉丝：\n• 极致沉浸感：理解球员视野、决策速度与空间感\n• 情感连接：与偶像同场战斗，大幅提升黏性\n• 新手友好：快速懂球\n\n对赛事方：\n• 票务溢价（AR体验专区）\n• 数据采集与精准营销\n\n参与度预估提升15-30%  |  AR体验票溢价40-60%"
    },
    7: {  # 幻灯片8
        "title": "亮点与创新点",
        "body": "三大亮点：\n01 从看球到代入球：史上首次大规模第一人称实时观赛\n02 混合现实完美融合：真实赛场+数字POV无缝切换\n03 个性化与社交：粉丝可分享POV高光片段\n\n差异化竞争力：超越传统转播和手机AR，制造现场独有稀缺感"
    },
    8: {  # 幻灯片9
        "title": "潜在挑战与解决方案",
        "body": "设备成本与舒适度 → 租赁模式 + 轻量化眼镜\n信息过载与晕动 → 可自定义显示层级 + AI智能过滤\n球员隐私与接受度 → AI重建为主，非真实头戴摄像头\n大规模并发处理 → 分阶段试点（VIP→全场）+ 边缘计算"
    },
    9: {  # 幻灯片10
        "title": "预期效果与商业价值",
        "body": "预期效果：观众停留时间↑25%、满意度↑35%、分享率↑50%\n\n商业模式：\n• 门票增值（AR专区溢价）\n• 付费内容订阅（球星POV套餐）\n• 品牌AR植入 + 虚拟商品\n\n未来展望：2026 NBA POV Mode首秀 → 2027顶级联赛常规配置 → 2028+全球体育赛事标配"
    },
    10: {  # 幻灯片11
        "title": "总结与建议",
        "body": "AR眼镜第一人称POV观赛是体育赛事从观看走向沉浸共创的关键一步，已在2026 NBA技术上落地验证。\n\n行动建议：\n• 立即启动VIP试点项目，选择2-3场本地赛事\n• 与技术伙伴（Meta/Xreal/5G运营商）对接测试\n• 结合本地篮球赛事快速验证并收集反馈数据\n• 建立商业合作框架：赞助+票务+内容付费"
    },
    11: {  # 幻灯片12: 结束页
        "texts": [
            ("汇报人：", "AR POV 体育创新团队"),
            ("汇报时间：", "2026.06"),
            ("汇报结束，感谢观看", "期待AR篮球POV早日落地"),
        ],
        "special": "thanks_page"
    },
}

# 处理每一张幻灯片
for slide_idx, slide_data in content.items():
    if slide_idx >= len(slides):
        print(f"  [SKIP] Slide {slide_idx+1} not in template (only {len(slides)} slides)")
        continue

    slide = slides[slide_idx]

    if slide_data.get("special") == "title_page":
        # 标题页：找到"汇报人："后面的占位
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for old, new in slide_data.get("texts", []):
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                print(f"  Slide{slide_idx+1}: '{old}' -> '{new}'")

    elif slide_data.get("special") == "thanks_page":
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for old, new in slide_data.get("texts", []):
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                print(f"  Slide{slide_idx+1}: '{old}' -> '{new}'")

    else:
        # 内容页：找标题占位符和正文占位符
        title = slide_data.get("title", "")
        body = slide_data.get("body", "")

        # 收集所有含文本的形状
        text_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                full_text = shape.text_frame.text.strip()
                if full_text:
                    text_shapes.append((shape, full_text))

        # 策略：找到标题占位符（含"请输入标题"或"单击添加标题"）替换为标题
        # 找到正文占位符（最长的那段占位文本）替换为正文
        title_replaced = False
        body_replaced = False

        for shape, full_text in text_shapes:
            if not title_replaced and ("请输入标题" in full_text or "单击添加标题" in full_text):
                # 替换标题
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        if "请输入标题" in run.text or "单击添加标题" in run.text:
                            run.text = title
                            title_replaced = True
                            print(f"  Slide{slide_idx+1}: Title -> '{title}'")
                            break
                    if title_replaced:
                        break

            elif not body_replaced and len(full_text) > 30 and ("请输入" in full_text or "根据实际情况" in full_text or "为确保演示效果" in full_text):
                # 替换正文
                tf = shape.text_frame
                # 保留第一个段落，替换其内容
                if tf.paragraphs:
                    first_para = tf.paragraphs[0]
                    if first_para.runs:
                        first_para.runs[0].text = body
                        # 清空后续 runs
                        for run in first_para.runs[1:]:
                            run.text = ""
                        body_replaced = True
                        print(f"  Slide{slide_idx+1}: Body replaced ({len(body)} chars)")

        # 如果还有"请输入标题"没替换（模板中可能用小标题占位符）
        if not title_replaced and title:
            for shape, full_text in text_shapes:
                if "请输入标题" in full_text or "单击添加标题" in full_text:
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if "请输入标题" in run.text or "单击添加标题" in run.text:
                                run.text = title
                                title_replaced = True
                                print(f"  Slide{slide_idx+1}: Title (2nd pass) -> '{title}'")
                                break
                        if title_replaced:
                            break
                    if title_replaced:
                        break

        # 如果有"关键词"文本，清空它
        for shape, full_text in text_shapes:
            if "关键词" in full_text:
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        if "关键词" in run.text:
                            run.text = run.text.replace("关键词", "")
                            print(f"  Slide{slide_idx+1}: Removed '关键词'")

    # 处理通用文本替换
    if "texts" in slide_data and "special" not in slide_data:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for old, new in slide_data.get("texts", []):
                            if old in run.text:
                                run.text = run.text.replace(old, new)

# 保存
prs.save(OUTPUT)
print(f"\nDone! Saved to: {OUTPUT}")
